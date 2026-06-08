from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colores
AZUL_OSCURO   = RGBColor(0x1E, 0x3A, 0x5F)
AZUL_MEDIO    = RGBColor(0x2D, 0x6A, 0x4F)
AZUL_CLARO    = RGBColor(0x2B, 0x6C, 0xB0)
BLANCO        = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLARO    = RGBColor(0xF1, 0xF5, 0xF9)
VERDE         = RGBColor(0x16, 0xA3, 0x4A)
NARANJA       = RGBColor(0xEA, 0x58, 0x0C)
GRIS_TEXTO    = RGBColor(0x47, 0x55, 0x69)
AZUL_ACENTO   = RGBColor(0x38, 0xBD, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # layout en blanco

# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=BLANCO,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_bullet_box(slide, items, l, t, w, h, size=16, color=GRIS_TEXTO, icon="•"):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color

def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.4, AZUL_OSCURO)
    add_rect(slide, 0, 1.4, 13.33, 0.07, AZUL_ACENTO)
    add_text(slide, title, 0.4, 0.12, 12, 0.8, size=30, bold=True,
             color=BLANCO, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12, 0.5, size=14,
                 color=AZUL_ACENTO, align=PP_ALIGN.LEFT)

def footer(slide, tiempo=""):
    add_rect(slide, 0, 7.1, 13.33, 0.4, AZUL_OSCURO)
    add_text(slide, "Arriendos360  |  Universidad Distrital  |  Ingeniería de Software I  |  2026",
             0.3, 7.12, 10, 0.3, size=10, color=RGBColor(0xAA,0xC4,0xE8), align=PP_ALIGN.LEFT)
    if tiempo:
        add_text(slide, tiempo, 11.5, 7.12, 1.5, 0.3, size=10,
                 color=AZUL_ACENTO, align=PP_ALIGN.RIGHT)

def tag(slide, text, l, t, color=AZUL_CLARO):
    add_rect(slide, l, t, len(text)*0.115 + 0.3, 0.38, color)
    add_text(slide, text, l+0.1, t+0.04, len(text)*0.115+0.1, 0.3,
             size=11, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# SLIDE 1 — PORTADA
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, AZUL_OSCURO)
add_rect(s, 0, 4.8, 13.33, 2.7, RGBColor(0x17, 0x2A, 0x4A))
add_rect(s, 0.5, 3.1, 0.08, 1.5, AZUL_ACENTO)

add_text(s, "ARRIENDOS360", 0.8, 1.8, 12, 1.1, size=52, bold=True,
         color=BLANCO, align=PP_ALIGN.LEFT)
add_text(s, "Gestión Inteligente de Arrendamientos", 0.8, 2.85, 10, 0.6,
         size=22, color=AZUL_ACENTO, align=PP_ALIGN.LEFT)

add_text(s, "Laura Alejandra Merchán León\nJuan Sebastián Díaz Rodríguez\nMarlon Alexander Herrera Choachi",
         0.8, 4.95, 8, 1.2, size=15, color=RGBColor(0xBF,0xDB,0xF7), align=PP_ALIGN.LEFT)
add_text(s, "Ingeniería de Software I  |  Universidad Distrital  |  Junio 2026",
         0.8, 6.2, 10, 0.5, size=13, color=RGBColor(0x7A,0xA8,0xD1), align=PP_ALIGN.LEFT)

for i, (x, op) in enumerate([(10.5,0.15),(11.3,0.1),(12.1,0.07),(11.0,0.2),(12.5,0.12)]):
    c = RGBColor(0x2B+i*10, 0x6C+i*5, 0xB0+i*5)
    add_rect(s, x, 0.3+i*0.8, 0.6, 0.6, c)

# ─────────────────────────────────────────────
# SLIDE 2 — EL PROBLEMA
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "El Problema", "¿Por qué existe Arriendos360?")
footer(s, "0:20 – 0:50")

cards = [
    ("📋", "Gestión manual", "Propietarios usan Excel y\nWhatsApp para controlar\nsus arriendos"),
    ("🏦", "Costos elevados", "Inmobiliarias cobran 8-10%\nmensual del canon de\narrendamiento"),
    ("📄", "Contratos físicos", "Documentos que se pierden,\nalteran o simplemente\nno se cumplen"),
]
for i, (icon, title, body) in enumerate(cards):
    x = 0.5 + i * 4.2
    add_rect(s, x, 1.7, 3.8, 4.5, BLANCO)
    add_rect(s, x, 1.7, 3.8, 0.07, AZUL_CLARO)
    add_text(s, icon, x+1.5, 1.9, 1, 0.7, size=32, color=AZUL_CLARO, align=PP_ALIGN.CENTER)
    add_text(s, title, x+0.2, 2.7, 3.4, 0.5, size=18, bold=True,
             color=AZUL_OSCURO, align=PP_ALIGN.CENTER)
    add_text(s, body, x+0.2, 3.3, 3.4, 1.5, size=14, color=GRIS_TEXTO,
             align=PP_ALIGN.CENTER)

add_rect(s, 1.5, 6.3, 10.3, 0.55, AZUL_CLARO)
add_text(s, "💡  Arriendos360: plataforma web que digitaliza contratos, pagos y estadísticas sin costos de intermediación",
         1.7, 6.35, 10, 0.45, size=13, bold=True, color=BLANCO, align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────
# SLIDE 3 — PMP
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Plan de Gestión (PMP)", "Modelo iterativo incremental · 14 semanas · 420 horas")
footer(s, "0:50 – 1:50")

# Tabla sprints
headers = ["Sprint", "Semanas", "Módulo", "Story Points", "Responsable"]
rows = [
    ["Sprint 1", "8 – 9", "Contratos + Auth", "11 pts", "Laura + Sebastián"],
    ["Sprint 2", "10 – 11", "Módulo Financiero", "8 pts", "Sebastián"],
    ["Sprint 3", "12", "Dashboard + Infra", "5 pts", "Marlon + Laura"],
    ["Total", "5 sem.", "MVP completo", "24 pts", "Todo el equipo"],
]
col_w = [1.4, 1.4, 2.8, 1.8, 2.8]
col_x = [0.5, 1.9, 3.3, 6.1, 7.9]
row_h, row_y0 = 0.52, 1.7

for ci, (h, x, w) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s, x, row_y0, w-0.05, row_h, AZUL_OSCURO)
    add_text(s, h, x+0.08, row_y0+0.08, w-0.2, row_h-0.1,
             size=12, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    y = row_y0 + (ri+1)*row_h
    bg = BLANCO if ri % 2 == 0 else RGBColor(0xE8,0xF0,0xFB)
    if ri == 3: bg = RGBColor(0xD1,0xFA,0xE5)
    for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        add_rect(s, x, y, w-0.05, row_h, bg)
        c = AZUL_OSCURO if ri == 3 else GRIS_TEXTO
        add_text(s, cell, x+0.08, y+0.1, w-0.2, row_h-0.15,
                 size=12, bold=(ri==3), color=c, align=PP_ALIGN.CENTER)

# Cards resumen
kpis = [("420 hrs", "Esfuerzo total"), ("3", "Integrantes"), ("14", "Semanas"), ("$0", "Presupuesto efectivo")]
for i, (val, lbl) in enumerate(kpis):
    x = 0.5 + i*3.2
    add_rect(s, x, 4.8, 3.0, 1.3, BLANCO)
    add_rect(s, x, 4.8, 3.0, 0.06, AZUL_ACENTO)
    add_text(s, val, x+0.1, 4.9, 2.8, 0.7, size=28, bold=True,
             color=AZUL_CLARO, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+0.1, 5.6, 2.8, 0.4, size=11,
             color=GRIS_TEXTO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 4 — SRS
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Especificación de Requisitos (SRS)", "19 Requisitos Funcionales · Priorización MoSCoW")
footer(s, "1:50 – 2:40")

must = ["RF-01: Registro de propietario","RF-03: Inicio de sesión JWT","RF-05: Roles diferenciados",
        "RF-06: Registro de contratos","RF-09: Estados del contrato",
        "RF-10: Registro de pagos","RF-12: Motor de mora","RF-14: Dashboard KPIs","RF-17: CRUD inmuebles"]
should = ["RF-07: Carga PDF contrato","RF-11: Recibo de pago","RF-08: Validación PDF"]
could  = ["RF-15: Gráficos dinámicos","RF-16: Alertas visuales con color"]

add_rect(s, 0.4, 1.6, 5.8, 5.0, BLANCO)
add_rect(s, 0.4, 1.6, 5.8, 0.45, VERDE)
add_text(s, "✅  MUST HAVE (Obligatorio)", 0.55, 1.65, 5.5, 0.35,
         size=13, bold=True, color=BLANCO)
add_bullet_box(s, must, 0.6, 2.15, 5.4, 4.2, size=12, color=GRIS_TEXTO, icon="✔")

add_rect(s, 6.7, 1.6, 5.8, 2.3, BLANCO)
add_rect(s, 6.7, 1.6, 5.8, 0.45, AZUL_CLARO)
add_text(s, "🔵  SHOULD HAVE (Deseable)", 6.85, 1.65, 5.5, 0.35,
         size=13, bold=True, color=BLANCO)
add_bullet_box(s, should, 6.9, 2.15, 5.4, 1.6, size=12, color=GRIS_TEXTO, icon="◎")

add_rect(s, 6.7, 4.1, 5.8, 1.9, BLANCO)
add_rect(s, 6.7, 4.1, 5.8, 0.45, NARANJA)
add_text(s, "🟠  COULD HAVE (Opcional)", 6.85, 4.15, 5.5, 0.35,
         size=13, bold=True, color=BLANCO)
add_bullet_box(s, could, 6.9, 4.65, 5.4, 1.2, size=12, color=GRIS_TEXTO, icon="◌")

add_rect(s, 6.7, 6.15, 5.8, 0.95, RGBColor(0xFE,0xE2,0xE2))
add_rect(s, 6.7, 6.15, 5.8, 0.06, RGBColor(0xEF,0x44,0x44))
add_text(s, "🔴  WON'T HAVE: Pasarela de pagos, DIAN, SMS, App móvil",
         6.85, 6.22, 5.5, 0.8, size=11, color=RGBColor(0x7F,0x1D,0x1D))

# ─────────────────────────────────────────────
# SLIDE 5 — SDD
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Diseño del Software (SDD)", "Arquitectura de 3 capas · Docker · PostgreSQL")
footer(s, "2:40 – 3:20")

capas = [
    ("CAPA DE PRESENTACIÓN", "React 18 · React Router · Axios", AZUL_CLARO,
     ["Login.js", "Dashboard.js", "Contratos.js", "Pagos.js", "Inmuebles.js"]),
    ("CAPA DE LÓGICA", "Node.js · Express · JWT · Multer", RGBColor(0x7C,0x3A,0xED),
     ["auth.controller", "contrato.controller", "pago.controller", "dashboard.controller"]),
    ("CAPA DE DATOS", "PostgreSQL · Sequelize ORM", RGBColor(0x0E,0x7A,0x90),
     ["usuarios", "contratos", "pagos", "inmuebles", "propietarios"]),
]
for i, (title, tech, color, items) in enumerate(capas):
    x = 0.4 + i * 4.3
    add_rect(s, x, 1.6, 4.0, 4.8, BLANCO)
    add_rect(s, x, 1.6, 4.0, 0.5, color)
    add_text(s, title, x+0.1, 1.65, 3.8, 0.35, size=11, bold=True,
             color=BLANCO, align=PP_ALIGN.CENTER)
    add_text(s, tech, x+0.1, 2.15, 3.8, 0.4, size=11,
             color=color, align=PP_ALIGN.CENTER, italic=True)
    add_bullet_box(s, items, x+0.2, 2.65, 3.6, 3.5, size=12, color=GRIS_TEXTO, icon="▸")
    if i < 2:
        add_text(s, "→", x+3.85, 3.6, 0.5, 0.5, size=22, bold=True,
                 color=AZUL_CLARO, align=PP_ALIGN.CENTER)

add_rect(s, 1.0, 6.5, 11.3, 0.5, AZUL_OSCURO)
add_text(s, "🐳  Docker Compose  ·  arriendos360_db  ·  arriendos360_api  ·  arriendos360_ui",
         1.2, 6.55, 11.0, 0.4, size=13, bold=True, color=BLANCO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 6 — DEMO
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, AZUL_OSCURO)
add_rect(s, 0, 0, 13.33, 7.5, RGBColor(0x1E,0x3A,0x5F))

add_text(s, "▶", 5.5, 2.2, 2.5, 2.0, size=80, color=AZUL_ACENTO, align=PP_ALIGN.CENTER)
add_text(s, "DEMO EN VIVO", 0, 4.4, 13.33, 0.9, size=38, bold=True,
         color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, "Navegación por la aplicación", 0, 5.25, 13.33, 0.5, size=18,
         color=AZUL_ACENTO, align=PP_ALIGN.CENTER)

steps = ["Login", "Dashboard", "Inmueble", "Contrato + PDF", "Cobro", "Pago", "Mora", "Recibo"]
for i, step in enumerate(steps):
    x = 0.5 + (i % 4) * 3.1
    y = 1.5 if i < 4 else 2.05
    tag(s, f"{i+1}. {step}", x, y, AZUL_CLARO if i % 2 == 0 else RGBColor(0x0E,0x7A,0x90))

footer(s, "3:20 – 7:20")

# ─────────────────────────────────────────────
# SLIDE 7 — TRAZABILIDAD
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Trazabilidad", "PMP → SRS → SDD → Código")
footer(s, "7:20 – 8:00")

cols  = ["PMP (Sprint)", "SRS (Requisito)", "SDD (Componente)", "Código"]
data  = [
    ["Sprint 1\nSem 8-9", "RF-01, RF-03, RF-05\nRF-06, RF-07, RF-09", "Auth Service\nContract Service", "auth.controller.js\ncontrato.controller.js"],
    ["Sprint 2\nSem 10-11", "RF-10, RF-12, RF-13\nRF-17, RF-18", "Financial Service\nInmuebles Service", "pago.controller.js\ninmueble.controller.js"],
    ["Sprint 3\nSem 12", "RF-14, RF-16", "Dashboard Service", "dashboard.controller.js"],
]
col_colors = [AZUL_OSCURO, AZUL_CLARO, RGBColor(0x7C,0x3A,0xED), RGBColor(0x0E,0x7A,0x90)]
cw = [2.5, 3.2, 3.2, 3.2]
cx = [0.35, 2.95, 6.25, 9.55]

for ci, (h, x, w, c) in enumerate(zip(cols, cx, cw, col_colors)):
    add_rect(s, x, 1.6, w-0.05, 0.55, c)
    add_text(s, h, x+0.1, 1.65, w-0.2, 0.42, size=13, bold=True,
             color=BLANCO, align=PP_ALIGN.CENTER)

for ri, row in enumerate(data):
    y = 2.25 + ri * 1.4
    bg = BLANCO if ri % 2 == 0 else RGBColor(0xEF,0xF6,0xFF)
    for ci, (cell, x, w) in enumerate(zip(row, cx, cw)):
        add_rect(s, x, y, w-0.05, 1.3, bg)
        add_text(s, cell, x+0.1, y+0.15, w-0.2, 1.1, size=11,
                 color=GRIS_TEXTO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 8 — RESULTADOS
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Resultados", "Lo que logramos en 14 semanas")
footer(s, "8:00 – 8:20")

kpis = [
    ("12/19", "Requisitos\nfuncionales\ncompletados", VERDE),
    ("83%", "Story Points\nentregados\n(20 de 24)", AZUL_CLARO),
    ("15+", "Endpoints\ndocumentados\nen Postman", RGBColor(0x7C,0x3A,0xED)),
    ("3", "Contenedores\nDocker\nfuncionando", RGBColor(0x0E,0x7A,0x90)),
]
for i, (val, lbl, color) in enumerate(kpis):
    x = 0.5 + i * 3.15
    add_rect(s, x, 1.7, 3.0, 4.5, BLANCO)
    add_rect(s, x, 1.7, 3.0, 0.08, color)
    add_text(s, val, x+0.1, 2.1, 2.8, 1.3, size=48, bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s, lbl, x+0.1, 3.5, 2.8, 1.5, size=14,
             color=GRIS_TEXTO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 9 — LECCIONES APRENDIDAS
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, GRIS_CLARO)
header_bar(s, "Lecciones Aprendidas", "Start · Stop · Continue")
footer(s, "8:20 – 9:00")

lecciones = [
    ("Laura", "🗄️", AZUL_CLARO,
     "Docker facilita el entorno compartido, pero es clave configurar las redes de contenedores desde el inicio del proyecto."),
    ("Sebastián", "🧪", RGBColor(0x7C,0x3A,0xED),
     "Los tests automatizados deben integrarse al sprint de QA, no dejarse en rama separada hasta el cierre."),
    ("Marlon", "📊", VERDE,
     "Los gráficos del frontend requieren más tiempo del estimado. Deben tener su propio Story Point desde la planeación."),
]
for i, (nombre, icon, color, texto) in enumerate(lecciones):
    x = 0.4 + i * 4.25
    add_rect(s, x, 1.6, 4.0, 4.9, BLANCO)
    add_rect(s, x, 1.6, 4.0, 0.08, color)
    add_text(s, icon, x+1.5, 1.75, 1.2, 0.9, size=36,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s, nombre, x+0.1, 2.75, 3.8, 0.5, size=18, bold=True,
             color=AZUL_OSCURO, align=PP_ALIGN.CENTER)
    add_text(s, texto, x+0.2, 3.35, 3.6, 2.8, size=13,
             color=GRIS_TEXTO, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# SLIDE 10 — CIERRE
# ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, AZUL_OSCURO)
add_rect(s, 0, 5.5, 13.33, 2.0, RGBColor(0x17,0x2A,0x4A))
add_rect(s, 0.5, 3.2, 0.08, 1.6, AZUL_ACENTO)

add_text(s, "¡Gracias!", 0, 1.2, 13.33, 1.4, size=58, bold=True,
         color=BLANCO, align=PP_ALIGN.CENTER)
add_text(s, "Arriendos360  ·  Junio 2026", 0, 2.65, 13.33, 0.6,
         size=18, color=AZUL_ACENTO, align=PP_ALIGN.CENTER)

stack = ["React 18", "Node.js + Express", "PostgreSQL", "Docker Compose", "JWT + bcrypt"]
for i, tech in enumerate(stack):
    x = 1.5 + i * 2.1
    add_rect(s, x, 3.7, 1.9, 0.45, RGBColor(0x2D,0x4E,0x7A))
    add_text(s, tech, x+0.05, 3.74, 1.8, 0.37, size=11, bold=True,
             color=AZUL_ACENTO, align=PP_ALIGN.CENTER)

add_rect(s, 2.0, 5.7, 9.3, 0.6, RGBColor(0x1E,0x3A,0x5F))
add_text(s, "🔗  github.com/UDFJDC-IngenieriaSoftware/IS_202610_E2",
         2.2, 5.75, 9.0, 0.5, size=14, color=AZUL_ACENTO, align=PP_ALIGN.CENTER)

add_text(s, "Laura Merchán  ·  Sebastián Díaz  ·  Marlon Herrera",
         0, 6.55, 13.33, 0.5, size=13, color=RGBColor(0x7A,0xA8,0xD1), align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────
# GUARDAR
# ─────────────────────────────────────────────
out = r'c:\Users\laura\OneDrive\Escritorio\Ingenieria de Software\IS_202610_E2\Presentacion-Arriendos360.pptx'
prs.save(out)
print(f"OK - Presentacion guardada en:\n{out}")
